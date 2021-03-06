from __future__ import print_function
import os, fnmatch, pptx, glob, tkinter, tempfile, shutil
import tkinter as tk, pandas as pd
import xml.etree.ElementTree as ET
from PyPDF2 import PdfFileReader
from pptx import Presentation
from pdf2image import convert_from_path
from tkinter import *
from tkinter import filedialog
from tkinter import *
from PyPDF2 import PdfFileReader
from tkinter import filedialog
from tkinter import ttk
import glob
import tqdm


from pptx import Presentation
import argparse
from tkinter import filedialog
import os








def analyze_ppt(dir):

    for root, dirs, files in os.walk(dir):
        for f in files:
            if f.endswith("pptx"):
                """ Take the input file and analyze the structure."""
                prs = Presentation()
                print(prs)
                # Each powerpoint file has multiple layouts
                # Loop through them all and  see where the various elements are
                for index, _ in enumerate(prs.slide_layouts):
                    slide = prs.slides.add_slide(prs.slide_layouts[index])
                    # Not every slide has to have a title
                    try:
                        title = slide.shapes.title
                        title.text = 'Title for Layout {}'.format(index)
                    except AttributeError:
                        print("No Title for Layout {}".format(index))
                    # Go through all the placeholders and identify them by index and type
                    for shape in slide.placeholders:
                        if shape.is_placeholder:
                            phf = shape.placeholder_format
                            # Do not overwrite the title which is just a special placeholder
                            try:
                                if 'Title' not in shape.text:
                                    shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                            except AttributeError:
                                print("{} has no text attribute".format(phf.type))
                            print('{} {}'.format(phf.idx, shape.name))
               # prs.save(output)




'''
# for ImportZip
import zipfile
from StringIO import StringIO
from PIL import Image
import imghdr
'''


'''
#Modules for MST import
from office365.runtime.auth.authentication_context import AuthenticationContext
from office365.sharepoint.client_context import ClientContext
from office365.sharepoint.file import File

# Connect to Microsoft Teams 365
client_id = 'yourclientid'
client_secret = 'yourclientsecret'
url = 'https://assettellc-my.sharepoint.com/teams/......yourteam'
relative_url = '/teams/yourteam/Shared%20Documents/yourteamschannel/yourdoc.extension'

#Ports file to local drive
ctx_auth = AuthenticationContext(url)
if ctx_auth.acquire_token_for_app(client_id, client_secret):
    ctx = ClientContext(url, ctx_auth)
    with open(filename, 'wb') as output_file:
        response = File.open_binary(ctx, relative_url)
        output_file.write(response.content)
else:
    print(ctx_auth.get_last_error())
'''
import glob
import tqdm
import win32com.client
import pywintypes
from os import getcwd


def main():
    dir = filedialog.askdirectory()

# Count total files by type in root and all sub-directories
    #Todo:
        #Group by Sub-directories
        #Speed up file counting
            #list_of_pdf_filenames = glob.glob('*.*') ???
        #Supress filetypes that were not found in Dataframe
    df = pd.DataFrame(columns=['FileType', 'Number'])

    # File format is case sensitive
    list = ["pdf", "ppt", "pptx", "jpg", "jpeg", "png", "svg", 'tif', "xml",
            "xlsm", "zip", "7z", "arj", "deb", "pkg", "rar", "csv", "dat",
            "log", "mdb", "sav", "sql", "tar", "apk", "bat", "bin", "cgi",
            "com", "exe", "jar", "msi", "py", "txt"]
    for l in list:

        i = 0
        for root, dirs, files in os.walk(dir):
            for file in files:
                if file.endswith(l):
                    i += 1

        df2 = pd.DataFrame([[l, i]],
                           columns=['FileType', 'Number'])
        df = df.append(df2, ignore_index=True)
    print(f'{df} \n  ----File Count Complete-----')

# Returns table of all PDFs with page number of each
# Returns empty Dataframe if no PDFs found
    list_of_pdf_filenames = glob.glob('*pdf')
    df = pd.DataFrame(columns=['fileName', 'fileLocation', 'pageNumber'])
    # PDF Data Frame
    for root, dirs, files in os.walk(dir):

        for f in files:
            if f.endswith(".pdf"):
                pdf = PdfFileReader(open(os.path.join(root, f), 'rb'))
                df2 = pd.DataFrame([[f, os.path.join(root, f), pdf.getNumPages()]],
                                   columns=['fileName', 'fileLocation', 'pageNumber'])
                df = df.append(df2, ignore_index=True)
    print(f'{df} \n  ----PDF Page Count Complete-----')



# Count PPT Elements by Keyword
    for root, dirs, files in os.walk(dir):
        for f in files:
            if f.endswith("pptx"):

                #prs = Presentation(f)
                #count = 0
                df = pd.DataFrame(columns=['File','Key', 'Count'])
               # f = os.path.basename(f)
                #keys are  case sensitive
                keys = ["Risk", "Performance", "Important", "Disclosure", "Portfolio", "Term 4", "Term 5, 7,8,9"]
                for k in keys:
                    count = 0

                    for slide in Presentation.slides:
                        title = slide.shapes.title.text
                        print(title)

                        if k in title:
                            count += 1

                    df2 = pd.DataFrame([[f,k, count]], columns=['File','Key', 'Count'])
                    df = df.append(df2, ignore_index=False)

                    if slide.shapes.title is None:
                        continue
                print(f'{df} \n  ----Process Complete-----')

'''
#Counts the number of pages in PPT 
# Not working correctly -- counts pages wrong
    df = pd.DataFrame(columns=['fileName', 'pageNumber'])
    ppt= 0
    for root, dirs, files in os.walk(dir):
        for f in files:
            if f.endswith(".pptx"):
                p = pptx.Presentation(os.path.abspath(f))
                ppt = len(p.slides)

                df2 = pd.DataFrame([[f, ppt]],
                                   columns=['fileName', 'pageNumber'])
                df = df.append(df2, ignore_index=True)
    print(f'{df} \n  ----PPT Page Count Complete-----')
'''
def PPTtoPDF():
#Works on UBUNTU 20.04 -- LibreOffice Dependent
    path = filedialog.askdirectory()
    files = [f for f in glob.glob(path + "**/*.pptx", recursive=True)]

    for f in tqdm.tqdm(files):
        command = "unoconv -f pdf \"{}\"".format(f)
        os.system(command)

    files = [f for f in glob.glob(path + "**/*.ppt", recursive=True)]
    for f in tqdm.tqdm(files):
        command = "unoconv -f pdf \"{}\"".format(f)
        os.system(command)
'''
def PPTtoPDF():
#Windows only
    formatType = 32
    inputFileName = filedialog.askopenfilename()
    powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
    powerpoint.Visible = 1
    outputFileName = os.path.basename(inputFileName)
    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

'''



def PPTtoTextParser():
    data = filedialog.askopenfilename()
    doc = ET.fromstring(data)
    for AAA in doc.findall('AAA'):
        print(len(AAA.findall('CCC')))








def CountPagesPDFWorking():
# Returns table of all PDFs with page number of each
    list_of_pdf_filenames = glob.glob('*pdf')
    df = pd.DataFrame(columns=['fileName', 'fileLocation', 'pageNumber'])

#PDF Data Frame
    for root, dirs, files in os.walk(filedialog.askdirectory()):

        for f in files:
            if f.endswith(".pdf"):

                pdf=PdfFileReader(open(os.path.join(root, f),'rb'))
                df2 = pd.DataFrame([[f, os.path.join(root,f), pdf.getNumPages()]], columns=['fileName', 'fileLocation', 'pageNumber'])
                df = df.append(df2, ignore_index=True)
    print(df.head)


    for root, dirs, files in os.walk(filedialog.askdirectory()):

        for f in files:
#PowerPoint Data Frame
# This works...retrieves the correct file and courn, but outputs the wrong number of pages on PPT
            if f.endswith(".pptx"):
                p = pptx.Presentation()
                df2 = pd.DataFrame([[f, os.path.join(root, f), len(p.slides)]],
                                   columns=['fileName', 'fileLocation', 'pageNumber'])
                df = df.append(df2, ignore_index=True)
    print(df.head)


def countSlides():
    #https://www.programmersought.com/article/20215432082/

    p = pptx.Presentation(filedialog.askopenfilename())

    print(len(p.slides))


def appendFileName(AppendDirectory):
    for root, dirs, files in os.walk(AppendDirectory):
        for f in files:
            #replaces '-' with a space, and appends the file with user input
            new_filename = f.replace("-", "")
            new_filename = input() + new_filename
            os.rename(f, new_filename)


def PDFpageExtraction():
    # this needs work: https://stackoverflow.com/questions/46184239/extract-a-page-from-a-pdf-as-a-jpeg
    # https://simply-python.com/2018/11/15/convert-pdf-pages-to-jpeg-with-python/


    filename = filedialog.askopenfilename()
    with tempfile.TemporaryDirectory() as path:
        images_from_path = convert_from_path(filename, output_folder=path, last_page=1, first_page=0)

    base_filename = os.path.splitext(os.path.basename(filename))[0] + '.jpg'

    save_dir = filedialog.askdirectory()

    for page in images_from_path:
        page.save(os.path.join(save_dir, base_filename), 'JPEG')

def extractZip():
    imgzip = open('100-Test.zip')
    zippedImgs = zipfile.ZipFile(imgzip)

    for i in xrange(len(zippedImgs.namelist())):
        print("iter", i, " ",)
        file_in_zip = zippedImgs.namelist()[i]
        if (".jpg" in file_in_zip or ".JPG" in file_in_zip):
            print("Found image: ", file_in_zip, " -- ",)
            data = zippedImgs.read(file_in_zip)
            dataEnc = StringIO(data)
            img = Image.open(dataEnc)
            print
            img
        else:
            print
            ""


def CompareImages():
    print("this doesn't work yet")
    folder_selected = filedialog.askdirectory()

class Bot(tk.Tk):

    def __init__(self, *args, **kwargs):
        tk.Tk.__init__(self, *args, **kwargs)
        container = tk.Frame(self)

        # Gets the requested values of the height and widht.
        windowWidth = self.winfo_reqwidth()
        windowHeight = self.winfo_reqheight()

        # Gets both half the screen width/height and window width/height-- the higher the divisor the more center(ish)
        positionRight = int(self.winfo_screenwidth() / 3 - windowWidth / 3)
        positionDown = int(self.winfo_screenheight() / 3 - windowHeight / 3)

        # Positions the window in the center of the page.
        self.geometry("+{}+{}".format(positionRight, positionDown))
        # This is the section of code which creates the main window
        container.pack(side="top", fill="both", expand=True)

        self.frames = {}

        for F in (StartPage, PageOne, Statistics,Configurations):
            frame = F(container, self)

            self.frames[F] = frame

            frame.grid(row=0, column=0, sticky="nsew")

        self.show_frame(StartPage)

    def show_frame(self, cont):
        frame = self.frames[cont]
        frame.tkraise()

class StartPage(tk.Frame):

    def __init__(self, parent, controller):

        tk.Frame.__init__(self, parent)

        x= os.getcwd()

        relPath = f'{x}\images\AssetteLogo.png'
        background_image = tk.PhotoImage(file = relPath)
        background_image = background_image.subsample(1, 1)
        background_label = tk.Label(self, image=background_image)
        background_label.place(x=0, y=0)
        background_label.image = background_image

        background_label.grid(row=0, column=0, padx=5, pady= 50, sticky ='n', rowspan = 1)

        def show():
            p = password.get()  # get password from entry
            e = email.get()

            if p == "" and e == "" and var1.get() == 1:
                controller.show_frame(PageOne)
            if p != "password1":
                print("wrong password")

        emailLabel = tk.Label(self, font=('arial', 16, 'normal'),text='Email:')
        emailLabel.grid(row=1, column=0, padx=115, pady=0, stick= "w")

        email = StringVar()
        email = Entry(self, font=('arial', 16, 'normal'), textvariable=email)
        email.grid(row=1, column=0, pady=10, padx=110,sticky ="e")


        passwordLabel = tk.Label(self, font=('arial', 16, 'normal'),text='Password:')
        passwordLabel.grid(row=2, column=0, padx=75, pady=10, stick= "w")
        password = StringVar()  # Password variable
        passwordEntry = Entry(self, font=('arial', 16, 'normal'), textvariable=password, show='*')
        passwordEntry.grid(row=2, column=0, pady=0, padx=110, sticky="e")



        link = tk.Label(self, text='I forgot My Password or Username.', fg = "royalblue3")
        link.bind("<1>", lambda event, text='I forgot My Password or Username.': self.click_link(event, text))
        link.grid(row=3, column=0, padx=170, pady= 0, sticky = "w")

        def print_selection():
            if (var1.get() == 1) & (term.get() == 0):
                print('Placeholder')
            elif (var1.get() == 0) & (term.get() == 1):
                print('Placeholder')
            elif (var1.get() == 1) & (term.get() == 1):
                print("Placeholder")
                editConfig()
            else:
                print('Placeholder')

        var1 = IntVar()
        term = IntVar()

        c1 = tk.Checkbutton(self, text='I have read and agree to the',
                            variable=var1, onvalue=1, offvalue=0, command=print_selection)
        c1.grid(row=4, column=0, padx=165, pady= 0, sticky ='w')

        checkbox2 = tk.Checkbutton(self, text='Remember my decision.',
                            variable=term, onvalue=1, offvalue=0, command=print_selection)

        checkbox2.grid(row=5, column=0, padx=165, pady= 0, sticky ='w')

        l1 = tk.Label(self, text='terms & conditions.', fg="royalblue3")
        l1.grid(row=4, column=0, padx=86, pady=10, sticky="e")
        l1.bind("<1>", lambda event, text='terms & conditions': self.terms_link(event, text))

        button = tk.Button(self, text="Log in", bg="lightgrey", width=20, command=show)
        button.grid(row=6, column=0, padx=10, pady=20)

    def click_link(self, event, text):
        print("You clicked '%s'" % text)
        #self.webdriver.Chrome("https://www.google.com")

    def terms_link(self, event, text):
        print("You clicked '%s'" % text)
        #self.webdriver.Chrome("https://www.google.com")

class PageOne(tk.Frame):

    def __init__(self, parent, controller):
        tk.Frame.__init__(self, parent)

        self.grid_columnconfigure(0, weight=1)

        x= os.getcwd()
        print(x)

        relPath = f'{x}\images\AssetteLogo.png'
        background_image = tk.PhotoImage(file = relPath)
        background_image = background_image.subsample(1, 1)
        background_label = tk.Label(self, image=background_image)
        #background_label.place(x=0, y=0)
        background_label.image = background_image

        background_label.grid(row=0, column=0, pady=50, padx=0, sticky ='', rowspan = 1)

    #    jobLabel = tk.Label(self, text='Job Title', font=('arial', 16, 'normal'))
     #   jobLabel.grid(row=1, column=0, pady=20, padx=0,sticky= "w")

        def getInfo():
          #  title = job_entry.get()
          #  location = job_location.get()
          #  salary = salary_entry.get()
          #  editConfig()
            print(typeVariable.get())
            print(so.current())

            print(subtypeVariable.get())

        def UploadAction(event=None):
            filename = filedialog.askopenfilename()
            print('Selected:', filename)

            f = (f"{os.getcwd()}\\resumes")
            print (f)

            shutil.copy(filename, f)

        uploadLabel = tk.Label(self, font=('arial', 16, 'normal'), text='Analyze:')
        uploadLabel.grid(row=1, column=0, padx=105, pady=10, stick="w")

        button = tk.Button(self, text='Select', command=UploadAction)
        button.grid(row=1, column=0, padx=98, pady=10, ipadx =98, stick="e")

        def typeCallBack(*args):
            if typeVariable.get() == 'Full Analysis':
                print(typeVariable.get())
                main()
            if typeVariable.get() == "Rapid Gap Analysis":
                print(typeVariable.get())
                if subtypeVariable.get() == "Count Files in Selected Directory":
                    print(subtypeVariable.get())
                    countFiles()
                if subtypeVariable.get() == "PPT Analysis (Not Working)":
                    print(subtypeVariable.get())
                    CountPagesPDFWorking()

            if typeVariable.get() == "Inspect File":

                if subtypeVariable.get() == "Analyze PPTX":
                    print(subtypeVariable.get())
                    dir = filedialog.askdirectory()
                   #input = filedialog.askopenfilename()
                    #output = os.path.basename(input[-4])
                    analyze_ppt(dir)#input, output)

                if subtypeVariable.get() == 'Append File Names':
                    print(subtypeVariable.get())

                if subtypeVariable.get() == 'Count Slides in PowerPoint':
                    print(subtypeVariable.get())
                    countSlides()
                if subtypeVariable.get() == "Compare Images":
                    print(subtypeVariable.get())
                    CompareImages()
                if subtypeVariable.get() == "PDF to Images":
                    print(subtypeVariable.get())
                    PDFpageExtraction()
                if subtypeVariable.get() == "PPT by Element":
                    print(subtypeVariable.get())
                    PPTbyElement()
                if subtypeVariable.get() == "Convert PDF to Text":
                    print(subtypeVariable.get())
                    #https://vitalflux.com/python-extract-text-pdf-file-using-pdfminer/
                    text = extract_text(filedialog.askopenfilename())
                    print(text)
                    #PPTtoTextParser()
                if subtypeVariable.get() == 'Count PPT Slides':
                    print(subtypeVariable.get())
                    countSlides()
                if subtypeVariable.get() == 'Convert PowerPoint to PDF (Ubuntu)':
                    print(subtypeVariable.get())
                    PPTtoPDF()





        typeSelectLabel = tk.Label(self, font=('arial', 16, 'normal'), text='Type:')
        typeSelectLabel.grid(row=2, column=0, padx=130, pady=10, stick="w")



        typeVariable = StringVar(self)
        typeVariable.set("Select From List")
        #typeVariable.trace('w', typeCallBack)

        so = ttk.Combobox(self, textvariable=typeVariable)
        so.config(values=('Rapid Gap Analysis', 'Full Analysis', 'Inspect File'))
        so.grid(row=2, column=0, padx=98, pady=10, ipadx = 48, stick="e")


        subtypeSelectLabel = tk.Label(self, font=('arial', 16, 'normal'), text='Sub-type:')
        subtypeSelectLabel.grid(row=3, column=0, padx=95, pady=10, stick="w")

        subtypeVariable = StringVar(self)
        subtypeVariable.set("Select From List")
       # subtypeVariable.trace('w', typeCallBack)

        so = ttk.Combobox(self, textvariable=subtypeVariable)
        so.config(values=('Count Files in Selected Directory',
                          'Analyze PPTX',
                          'PPT Analysis (Not Working)',
                          'Count Slides in PowerPoint',
                          'PDF Analysis',
                          'Append File Names',
                          'Compare Images (Not Working)',
                          'Covert PDF to Image (Not Working)',
                          'Covert PPT to Text (Not Working)',
                          'Convert PDF to Text',
                          'Convert PowerPoint to PDF (Ubuntu)', 'PPT by Element'))
        so.grid(row=3, column=0, padx=98, pady=10, ipadx =48, stick="e")


        AnalyzeButton = tk.Button(self,
                                text="Analyze",

                                command=typeCallBack)
        AnalyzeButton.configure(
                              fg="white", font = ('arial', 20, 'bold'),
                              bg = "royalblue3",
                              activebackground="royalblue2",
                              relief=FLAT)

        AnalyzeButton.grid(row=5, column=0, padx = 0, pady = 15)

        button1 = tk.Button(self, text="Statistics",
                            command=lambda: controller.show_frame(Statistics))
        button1.grid(row=6, column=0, ipadx =18,  padx =175, pady = 0, sticky = "w")

        button3 = tk.Button(self, text="Configurations",
                            command=lambda: controller.show_frame(Configurations))
        button3.grid(row=6, column=0, padx= 170, pady = 0, sticky= "e")


class Statistics(tk.Frame):

    def __init__(self, parent, controller):
        tk.Frame.__init__(self, parent, relief=RAISED, borderwidth=1)
        label = tk.Label(self, text="'Statistics'")
        label.grid(row=0, column=0, pady=10, padx=10)

        button2 = tk.Button(self, text="Main Page",
                            command=lambda: controller.show_frame(PageOne))
        button2.grid(row=2, column=0, padx = 10, pady = 10)









class Configurations(tk.Frame):


    def __init__(self, parent, controller):
        tk.Frame.__init__(self, parent)

        self.grid_columnconfigure(0, weight=1)

        x= os.getcwd()
        print(x)

        relPath = f'{x}\images\AssetteLogo.png'
        background_image = tk.PhotoImage(file = relPath)
        background_image = background_image.subsample(1, 1)
        background_label = tk.Label(self, image=background_image)
        #background_label.place(x=0, y=0)
        background_label.image = background_image


        background_label.grid(row=0, column=0, pady=50, padx=0, sticky ='', rowspan = 1)





        Blacklist = tk.Label(self, font=('arial', 16, 'normal'), text='Exclude:')
        Blacklist.grid(row=3, column=0, padx=99, pady=10, stick="w")

        blacklist = StringVar()
        salary_entry = Entry(self,
                      text='salary_entry',
                      font=('arial', 16, 'normal'), textvariable=blacklist)

        salary_entry.grid(row=3, column=0, pady=10, padx=100, sticky = "e")

        def saveInfo():
            print("saved")
            print("Domain value is:" + n.get())
            print(f"{blacklist.get()} was removed added to the 'do not apply to' list")

        ApplyButton = tk.Button(self,
                                text="Save",

                                command=saveInfo)
        ApplyButton.configure(
                              fg="white", font = ('arial', 20, 'bold'),
                              bg = "royalblue3",
                              activebackground="royalblue2",
                              relief=FLAT)
        #ApplyButton.pack()
        ApplyButton.grid(row=5, column=0, padx = 0, pady = 15)


        button1 = tk.Button(self, text="Statistics",
                            command=lambda: controller.show_frame(Statistics))
        button1.grid(row=6, column=0, ipadx =18,  padx =175, pady = 0, sticky = "w")


        button3 = tk.Button(self, text="Main Page",
                            command=lambda: controller.show_frame(PageOne))
        button3.grid(row=6, column=0, padx= 170, ipadx=13 , pady = 0, sticky= "e")





app = Bot()
app.mainloop()
